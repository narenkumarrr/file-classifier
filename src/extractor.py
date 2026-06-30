import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
from pytesseract import TesseractNotFoundError

MAX_CHARS = 2500

def extract_text(file_path: str) -> str:
    """Extracts text from various file formats, capped at MAX_CHARS."""
    if not os.path.exists(file_path):
        return ""

    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext in ['.txt', '.csv', '.md']:
            return extract_txt(file_path)
        elif ext == '.pdf':
            return extract_pdf(file_path)
        elif ext == '.docx':
            return extract_docx(file_path)
        elif ext == '.pptx':
            return extract_pptx(file_path)
        elif ext in ['.png', '.jpg', '.jpeg']:
            return extract_image(file_path)
        else:
            print(f"Unsupported file format: {ext}")
            return ""
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return ""

def extract_txt(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read(MAX_CHARS)

def extract_pdf(file_path: str) -> str:
    text = ""
    reader = PdfReader(file_path)
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
        if len(text) >= MAX_CHARS:
            break
    return text[:MAX_CHARS]

def extract_docx(file_path: str) -> str:
    text = ""
    doc = Document(file_path)
    for para in doc.paragraphs:
        if para.text.strip():
            text += para.text + "\n"
        if len(text) >= MAX_CHARS:
            break
    return text[:MAX_CHARS]

def extract_pptx(file_path: str) -> str:
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
        if len(text) >= MAX_CHARS:
            break
    return text[:MAX_CHARS]

def extract_image(file_path: str) -> str:
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return text[:MAX_CHARS]
    except TesseractNotFoundError:
        print("Warning: Tesseract OCR is not installed. Cannot extract text from image.")
        return ""
